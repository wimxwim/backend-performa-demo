#!/usr/bin/env python3
# presentasi/generate_pptx.py — Generate 35 slides Modul Performa Backend GR Demo
# Theme: navy #1e3a5f + orange #f59e0b, Calibri, 16:9, Title+Content
import pptx
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
NAVY_DARK = RGBColor(0x14, 0x26, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
ORANGE_LIGHT = RGBColor(0xFF, 0xFB, 0xEB)

prs = pptx.Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
# Use blank layout for full control
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    # full white bg
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()


def add_header(slide, title, subtitle=None):
    # navy top bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.95)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    # orange accent line
    acc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.95), prs.slide_width, Inches(0.07)
    )
    acc.fill.solid()
    acc.fill.fore_color.rgb = ORANGE
    acc.line.fill.background()
    # title
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.5), Inches(0.65))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    # slide number badge on header right
    # subtitle / bab label on header right
    if subtitle:
        tx2 = slide.shapes.add_textbox(
            Inches(10.0), Inches(0.28), Inches(3.0), Inches(0.4)
        )
        tf2 = tx2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.alignment = PP_ALIGN.RIGHT
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
        p2.font.name = "Calibri"


def add_footer(slide, bab, idx, total=35):
    # footer bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        prs.slide_height - Inches(0.32),
        prs.slide_width,
        Inches(0.32),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY_DARK
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(
        Inches(0.3), prs.slide_height - Inches(0.28), Inches(9.5), Inches(0.22)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Gotong Royong \u2014 OS Kehidupan Komunitas  |  60 menit  |  {bab}  |  TIGA INSAN: Muttaqin \u2022 Shalih \u2022 Nafi'"
    p.font.size = Pt(7)
    p.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p.font.name = "Calibri"
    # slide number
    tx2 = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.2),
        prs.slide_height - Inches(0.28),
        Inches(1.0),
        Inches(0.22),
    )
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f"{idx} / {total}"
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(7)
    p2.font.color.rgb = RGBColor(0xFD, 0xBA, 0x74)
    p2.font.name = "Calibri"
    p2.font.bold = True


def add_bullets(
    slide,
    bullets,
    left=Inches(0.5),
    top=Inches(1.35),
    width=Inches(7.2),
    height=Inches(5.4),
    font_size=Pt(11),
):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.font.size = font_size
        p.font.name = "Calibri"
        p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        p.level = 0
        # bullet char via paragraph
        p.text = "\u2022  " + b
    return tx


def add_card(slide, left, top, width, height, title, body, icon=""):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    shape.shadow.inherit_shadow = False
    # title
    tx = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.3)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (icon + " " if icon else "") + title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "Calibri"
    # body
    tx2 = slide.shapes.add_textbox(
        left + Inches(0.15),
        top + Inches(0.38),
        width - Inches(0.3),
        height - Inches(0.5),
    )
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"


def add_table_slide(slide, headers, rows, left, top, width, height, col_widths=None):
    tbl_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, width, height
    )
    tbl = tbl_shape.table
    # header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(7)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(7)
                    run.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
                    run.font.name = "Calibri"
            # zebra
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if col_widths:
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = w
    return tbl_shape


def add_mono_box(slide, text, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    shape.line.fill.background()
    tx = slide.shapes.add_textbox(
        left + Inches(0.12),
        top + Inches(0.08),
        width - Inches(0.24),
        height - Inches(0.16),
    )
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x86, 0xEF, 0xAC)


# ──────────────────────────────────────────────
# SLIDES DATA — 35 slides
# ──────────────────────────────────────────────
slides = []

# 1 Cover
slides.append(
    dict(
        title="Modul Performa Backend GR Demo",
        footer="Cover",
        cover=True,
        cover_title="Modul Performa\nBackend GR Demo",
        cover_sub="Logging + Performa  \u2022  TIGA INSAN  \u2022  60 menit",
        cover_bullets=[
            "Poster 1+2 + 10 Bab + 16 Endpoint + Checklist",
            "Postgres + Redis + pg_trgm + MatView + PgBouncer  \u2192  p50 <50ms",
            "CDC Debezium WAL \u2192 Kafka \u2192 ES geospasial + ClickHouse",
        ],
    )
)

# 2 Daftar Isi
slides.append(
    dict(
        title="Daftar Isi  \u2014  10 Bab + 35 Slides Map",
        footer="Daftar Isi",
        bullets=[
            "Bab 1  Kecepatan = Kepercayaan (Muttaqin, SLA)  \u2192  Slides 3-5",
            "Bab 2  Data Flow Flutter \u2192 Gateway \u2192 7 Fondasi \u2192 6 DB  \u2192  Slide 13",
            "Bab 3  Postgres Scale (B-Tree, pg_trgm, MatView, Cursor, RLS)  \u2192  Slides 14-17",
            "Bab 4  Caching (L1/L2/L3, TTL, tiering)  \u2192  Slide 18  \u2022  Bab 5 ES geospasial  \u2192  Slide 19",
            "Bab 6  CDC Streaming (WAL \u2192 Kafka \u2192 ES/CH, anti dual-write)  \u2192  Slide 20",
            "Bab 7  API Delivery (GZIP, cursor, Edge, rate limit)  \u2192  Slide 21  \u2022  Bab 8 Observability  \u2192  Slide 22",
            "Bab 9  Roadmap 5 Fase (MVP Rp0 \u2192 auto-scale)  \u2192  Slide 23  \u2022  Bab 10 SLA 16 endpoint + throughput + DoD  \u2192  Slides 24-26",
            "Demo 01-05 + Appendix A/B/C + Penutup  \u2192  Slides 27-35",
        ],
    )
)

# 3 Bab1.1
slides.append(
    dict(
        title="Bab 1.1  \u2014  Kecepatan = Kepercayaan (Muttaqin)",
        footer="Bab 1",
        bullets=[
            "Muttaqin = fondasi TIGA INSAN: beriman dengan akal yang hidup \u2014 kepercayaan harus lebih dulu dari fitur.",
            "SHA-256 hash chain transparansi keuangan: setiap transaksi kas diverifikasi <1 detik \u2192 jamaah percaya pengurus.",
            "Lambat = khianat amanah: setiap ms keterlambatan mengikis kepercayaan yang susah payah dibangun.",
            "Kecepatan adalah bahasa universal keandalan sistem \u2014 warga cek laporan kas, data muncul <1 detik = platform serius.",
            "Prinsip: ukur server-side (p50/p95/p99), bukan feeling \u2014 EXPLAIN ANALYZE, pg_stat_statements, prom-client.",
        ],
    )
)

# 4 Bab1.2
slides.append(
    dict(
        title="Bab 1.2  \u2014  Konteks Indonesia: 3G, 2-3GB RAM, Kuota",
        footer="Bab 1",
        bullets=[
            "Jaringan 3G RTT 500-1000ms, sinyal tidak stabil \u2014 backend harus <200ms agar total <3 detik (UX #46, #50).",
            "Android entry-level RAM 2-3GB, CPU terbatas \u2014 payload kecil + GZIP 70-80% = hemat kuota & parsing cepat.",
            "Kuota terbatas: setiap byte tidak perlu adalah pemborosan \u2014 payload shaping ?fields= + Brotli + Edge caching.",
            "Implikasi: jangan kirim SELECT *; jangan polling; cache jadwal sholat TTL 1 jam; kompresi wajib.",
            "Gotong Royong bukan Silicon Valley \u2014 performa adalah aksesibilitas, bukan kemewahan.",
        ],
    )
)

# 5 Bab1.3
slides.append(
    dict(
        title="Bab 1.3  \u2014  Target SLA: p50 / p95 / p99 / Availability",
        footer="Bab 1",
        bullets=[
            "Global: p50 <50ms (read) / <100ms (write), p95 <200ms, p99 <500ms \u2014 diukur server-side (tanpa RTT jaringan).",
            "Availability: MVP 99.5% (3.6 jam down/bulan) \u2192 Fase 3+ 99.9% (43 menit) \u2192 Fase 5 99.99% (4 menit).",
            "Error rate <0.1% semua endpoint; jika p95 >500ms 5 menit \u2192 incident, cek pg_stat_statements + EXPLAIN.",
            "Alat ukur: prom-client histogram (buckets 0.01-5s), Prometheus histogram_quantile, Grafana, k6/autocannon.",
            "16 endpoint punya SLA per-endpoint (Bab 10.1) \u2014 tidak ada satu angka untuk semua; baca tabel SLA.",
        ],
    )
)

# 6 Poster #1
slides.append(
    dict(
        title="Poster #1  \u2014  Apa itu 200ms",
        footer="Poster #1",
        bullets=[
            "1 detik = 1000ms  \u2014  200ms = 0.2 detik = kedipan mata.",
            "Diagram: User (tap) \u2192 Server (proses 50ms) \u2192 DB (query 20ms) \u2192 Response (GZIP 10ms) \u2192 User lihat hasil.",
            "200ms adalah batas 'terasa instan' \u2014 di atas itu pengguna mulai merasa menunggu (UX #46).",
            "Target backend <200ms p95 \u2192 sisa budget untuk jaringan 3G (500-1000ms) masih <3 detik total.",
            "Ukur dengan X-Response-Time header + prom-client histogram, bukan stopwatch manual.",
        ],
        mono="User --tap--> [Gateway JWT 5ms] --> [Redis HIT 2ms / DB 20ms] --> [GZIP 5ms] --> User\nTotal server: 30-50ms (p50)  |  + RTT 3G 500ms = ~550ms (masih <3s)",
    )
)

# 7 Poster #2
slides.append(
    dict(
        title="Poster #2  \u2014  Semakin kecil ms, semakin cepat",
        footer="Poster #2",
        bullets=[
            "Intuisi: angka ms kecil = cepat, besar = lambat \u2014 jangan tertukar.",
            "Tabel acuan (server-side): 1ms sangat cepat \u00b7 10ms cepat \u00b7 50ms baik \u00b7 200ms batas \u00b7 1000ms lambat \u00b7 10000ms sangat lambat.",
            "Contoh: cache HIT Redis 1-5ms (sangat cepat) vs Seq Scan 2000ms (sangat lambat) = 400-2000x beda.",
            "Goal: pindahkan sebanyak mungkin request dari kolom kanan (lambat) ke kiri (cepat) via index + cache.",
            "Grafana: garis p50 (hijau), p95 (kuning), p99 (merah) \u2014 kejar hijau, jaga kuning, waspada merah.",
        ],
        table_headers=["ms", "Rasa", "Contoh"],
        table_rows=[
            ["1-10 ms", "Sangat cepat", "Redis HIT, RLS <0.1ms"],
            ["10-50 ms", "Cepat", "GIN 10-50ms, MatView 5-30ms"],
            ["50-200 ms", "Baik", "p50 read, p95 target"],
            ["200-500 ms", "Mulai terasa", "p99 batas, perlu optimasi"],
            ["500-1000 ms", "Lambat", "OFFSET 2000ms, LIKE Seq Scan"],
            [">1000 ms", "Sangat lambat", "N+1, tanpa index"],
        ],
    )
)

# 8 Poster #3
slides.append(
    dict(
        title="Poster #3  \u2014  Apa itu P99",
        footer="Poster #3",
        bullets=[
            "100 request diurutkan dari tercepat ke terlambat \u2014 P99 = request ke-99 (99% lebih cepat, 1% lebih lambat).",
            "Contoh: P99 = 200ms artinya 99 request <=200ms, 1 request >200ms (yang paling lambat, paling sering komplain).",
            "Kenapa peduli P99? Karena 1% yang lambat adalah pengguna yang paling vokal \u2014 mereka yang churn.",
            "Jangan hanya lihat rata-rata (avg) \u2014 avg bisa 50ms tapi P99 2000ms = ada yang sangat menderita.",
            "Spec: p50 <50ms, p95 <200ms, p99 <500ms \u2014 ketiganya harus hijau, bukan cuma p50.",
        ],
        mono="100 request sorted: [10,12,15,...,45,48, 50(p50), ..., 180, 200(p99), 2500]\n  50% <= p50   95% <= p95   99% <= p99   1% tail latency (paling lambat)",
    )
)

# 9 Poster #4
slides.append(
    dict(
        title="Poster #4  \u2014  P50 / P95 / P99 / P99.9",
        footer="Poster #4",
        bullets=[
            "P50 (median): 50% request lebih cepat \u2014 gambaran 'mayoritas' pengguna.",
            "P95: 95% lebih cepat \u2014 5% masih lebih lambat; target p95 <200ms.",
            "P99: 99% lebih cepat \u2014 hanya 1% tail; target p99 <500ms.",
            "P99.9: 99.9% lebih cepat \u2014 hanya 0.1% (1 dari 1000) \u2014 untuk Fase 5 (200k RPS, tail sangat penting).",
            "Distribusi: mayoritas (p50) \u2192 banyak (p95) \u2192 sedikit (p99) \u2192 sangat sedikit (p99.9) \u2014 ekor panjang harus dipangkas.",
        ],
        mono="Distribusi latency (ms):\n p50=30  |███████ majority\n p95=150 |███\n p99=400 |█ tail\np99.9=900|· very tail  --> optimasi tail = index + cache + proteksi",
    )
)

# 10 Poster #5
slides.append(
    dict(
        title="Poster #5  \u2014  Berapa cepat / lambat (skala rasa)",
        footer="Poster #5",
        bullets=[
            "Skala rasa pengguna (server-side): 0-50ms sangat cepat, 50-100ms sangat baik, 100-200ms baik, 200-500ms mulai terasa, 500-1000ms lambat, >1000ms sangat lambat.",
            "Mapping ke SLA: p50 harus di 'sangat cepat/sangat baik', p95 di 'baik', p99 jangan sampai 'lambat'.",
            "Jika p95 >500ms \u2192 pengguna mulai komplain; >1000ms \u2192 tinggalkan aplikasi.",
            "Gunakan skala ini saat baca Grafana \u2014 warna hijau (0-100), kuning (100-200), oranye (200-500), merah (>500).",
            "Setiap ms yang dipangkas = kepercayaan yang ditambah (Muttaqin).",
        ],
        table_headers=["Rentang", "Rasa", "Aksi"],
        table_rows=[
            ["0-50 ms", "Sangat cepat", "Pertahankan (cache HIT, GIN, MatView)"],
            ["50-100 ms", "Sangat baik", "Ideal p50 read"],
            ["100-200 ms", "Baik", "Batas p95, masih OK"],
            ["200-500 ms", "Mulai terasa", "Optimasi: index, N+1, payload"],
            ["500-1000 ms", "Lambat", "Incident jika p95 di sini"],
            [">1000 ms", "Sangat lambat", "Wajib fix: Seq Scan, OFFSET"],
        ],
    )
)

# 11 Poster #6.1
slides.append(
    dict(
        title="Poster #6.1  \u2014  10 Metrik Wajib",
        footer="Poster #6",
        bullets=[
            "10 metrik Poster #6 adalah kompas performa \u2014 tanpa ukur, tidak bisa kelola.",
            "3 pilar observabilitas (Bab 8.5): Metrics (Prometheus), Logs (Loki), Traces (Jaeger) \u2014 saling melengkapi.",
            "Cache hit rate >80% adalah pembeda MVP yang scalable vs yang boros DB.",
            "DB Query Time p95 <50ms & no Seq Scan >1000 rows \u2014 cek pg_stat_statements tiap rilis.",
        ],
        table_headers=["#", "Metrik", "Target", "Alat"],
        table_rows=[
            ["1", "Response Time", "p50 <50ms read", "prom-client histogram"],
            ["2", "P95 / P99", "<200ms / <500ms", "histogram_quantile"],
            ["3", "Throughput RPS/QPS", "100 \u2192 200k", "http_requests_total, k6"],
            ["4", "Latency RTT/TTFB", "<200ms", "curl time_total"],
            ["5", "Error Rate", "<0.1%", "5xx / total"],
            ["6", "Availability", "99.5% \u2192 99.99%", "up, SLO"],
            ["7", "CPU Usage", "<70% avg", "node_cpu"],
            ["8", "Memory Usage", "<70% avg", "node_memory"],
            ["9", "DB Query Time", "<50ms p95", "pg_stat_statements"],
            ["10", "Cache Hit Rate", ">80%", "cache_hit_total"],
        ],
    )
)

# 12 Poster #7 Glossary
slides.append(
    dict(
        title="Poster #7  \u2014  Glossary Super Sederhana",
        footer="Poster #7",
        bullets=[
            "Bahasa performa tanpa jargon \u2014 untuk pengurus masjid & RT/RW juga paham.",
            "ms = milidetik (1/1000 detik); P99 = 99% request lebih cepat dari angka ini.",
            "RPS/QPS = request/detik; Latency = waktu bolak-balik; Throughput = berapa banyak yang bisa dilayani.",
            "Cache Hit = ambil dari memori cepat (tanpa tanya DB); Error Rate = % gagal; Availability = % waktu hidup.",
        ],
        table_headers=["Istilah", "Arti 1 kalimat"],
        table_rows=[
            ["ms", "Milidetik: 1000ms = 1 detik"],
            ["P99", "99% request <= angka ini, 1% lebih lambat"],
            ["RPS/QPS", "Request/Query per detik (throughput)"],
            ["Latency", "Waktu tunggu bolak-balik"],
            ["Throughput", "Kapasitas layani request/detik"],
            ["Error Rate", "% request gagal (5xx)"],
            ["Availability", "% waktu layanan hidup"],
            ["CPU / Memory", "% otak & ingatan server terpakai"],
            ["Cache Hit", "% dilayani dari memori cepat"],
        ],
    )
)

# 13 Bab2 Data Flow
slides.append(
    dict(
        title="Bab 2  \u2014  Data Flow: Flutter \u2192 Gateway \u2192 7 Fondasi \u2192 6 DB",
        footer="Bab 2",
        bullets=[
            "Flutter (single codebase) \u2192 API Gateway (Kong/Supabase Edge: JWT, rate limit, GZIP) \u2192 7 Fondasi Bersama.",
            "7 Fondasi: Auth, Profile, Payment (HarmoniPay+Xendit), Notification (FCM/Wablas), Storage (S3/R2), Audit SHA-256, Feature Flag.",
            "6 DB: Postgres (jantung ACID+RLS), Redis (<10ms cache), Mongo (kajian), ES (geo+search), ClickHouse (OLAP), Influx/Timescale (IoT).",
            "Bottleneck per layer: Gateway (JWT cache <5ms), Fondasi (profile cache 5m), DB (index+pool), Network (GZIP 70%).",
            "Hot path read: Flutter \u2192 Gateway cache \u2192 Redis HIT <5ms \u2192 response (tanpa DB).",
        ],
        mono="Flutter --> Gateway(JWT cache) --> [7 Fondasi] --> Redis HIT <5ms --> response\n                         \\--> Postgres(ACID) --WAL--> Debezium --> Kafka --> ES/ClickHouse (async)\n                         \\--> pg_trgm(MVP) / ES(Fase3) --> cache populer",
    )
)

# 14 Bab3.1 B-Tree
slides.append(
    dict(
        title="Bab 3.1  \u2014  Index B-Tree: 50.000x lebih cepat",
        footer="Bab 3",
        bullets=[
            "Tanpa index: Seq Scan baca 1.000.000 baris satu-per-satu \u2014 lambat, O(n).",
            "Dengan B-Tree: Balanced Tree cari dalam ~20 langkah (logaritmik) \u2014 O(log n) \u2192 50.000x speedup.",
            "Aturan: semua FK wajib punya index (Checklist #3) \u2014 tanpa index, JOIN = Seq Scan.",
            "Buat CONCURRENTLY agar tidak blokir write di produksi; cek EXPLAIN: Index Scan vs Seq Scan.",
            "Index komposit: (community_id, pinned, created_at) untuk ORDER BY pinned, created_at LIMIT 20 (p50 <30ms).",
        ],
        mono="Tanpa index: [1M rows] --scan--> 1.000.000 langkah ~2000ms\nDengan B-Tree:  root -> branch -> branch -> leaf  ~20 langkah ~10ms\nCREATE INDEX CONCURRENTLY idx_komunitas_slug ON communities(slug);",
    )
)

# 15 Bab3.2 pg_trgm
slides.append(
    dict(
        title="Bab 3.2  \u2014  pg_trgm GIN: 10-50ms vs LIKE 2000ms",
        footer="Bab 3",
        bullets=[
            "LIKE '%ayam%' tidak bisa pakai B-Tree \u2192 Seq Scan 6.081 baris ~2000ms (meledak saat data tumbuh).",
            "pg_trgm pecah teks jadi trigram (3 huruf), GIN index cari overlap trigram \u2192 10-50ms, similarity 71% untuk 'ayam'.",
            "Query: WHERE name % 'ayam' ORDER BY similarity(name,'ayam') DESC LIMIT 20 (operator % = similarity >0.3).",
            "Aktifkan: CREATE EXTENSION pg_trgm; CREATE INDEX USING GIN (name gin_trgm_ops) \u2014 gratis, built-in Postgres.",
            "Kapan ES? pg_trgm cukup untuk puluhan ribu dokumen; ES untuk 500+ komunitas / ratusan ribu dokumen (Fase 3).",
        ],
        mono="BEFORE: SELECT * FROM umkm WHERE name LIKE '%ayam%'  --> Seq Scan 2000ms\nAFTER : SELECT * FROM umkm WHERE name % 'ayam' ORDER BY similarity DESC --> GIN 10-50ms",
    )
)

# 16 Bab3.3-3.4
slides.append(
    dict(
        title="Bab 3.3-3.4  \u2014  MatView + Cursor vs OFFSET",
        footer="Bab 3",
        bullets=[
            "MatView: agregasi SUM(amount) GROUP BY community_id dari 500ms (Seq Scan) \u2192 5-30ms (Index Scan MatView) \u2014 refresh tiap 5 menit CONCURRENTLY.",
            "OFFSET 10000: scan 10.020 baris, buang 10.000 \u2192 2000ms; Cursor keyset: WHERE (created_at,id) > cursor LIMIT 20 \u2192 20ms (100x).",
            "Cursor di-encode base64: {created_at, id} \u2192 nextCursor; stabil, tidak skip/duplikat saat data baru masuk.",
            "Kapan OFFSET boleh? Hanya untuk halaman kecil (<100) atau export batch; untuk feed/list wajib cursor.",
            "Verifikasi: EXPLAIN (ANALYZE, BUFFERS) \u2014 cek Execution Time & Buffers hit.",
        ],
    )
)

# 17 Bab3.5-3.8
slides.append(
    dict(
        title="Bab 3.5-3.8  \u2014  PgBouncer, EXPLAIN, RLS, VACUUM",
        footer="Bab 3",
        bullets=[
            "PgBouncer pool 25, transaction mode :6432 \u2014 dari 200 koneksi langsung \u2192 25 pool, hemat memori & context switch.",
            "EXPLAIN ANALYZE: metode identifikasi query lambat \u2014 lihat Seq Scan, Index Scan, Execution Time, Buffers.",
            "RLS <0.1ms overhead: CREATE POLICY community_isolation USING (community_id = current_setting(...)) + index FK.",
            "VACUUM & ANALYZE: pemeliharaan rutin \u2014 VACUUM hapus dead tuple, ANALYZE update statistik planner agar tidak salah pilih Seq Scan.",
            "Checklist: autovacuum normal, n_dead_tup <1000, pg_stat_statements enabled, slow log >100ms.",
        ],
    )
)

# 18 Bab4 Caching
slides.append(
    dict(
        title="Bab 4  \u2014  Caching: Hierarki L1/L2/L3 + TTL + Tiering",
        footer="Bab 4",
        bullets=[
            "Hierarki: L1 memory sub-ms (hot), L2 Redis 1-5ms (warm), L3 Postgres 10-50ms (cold) \u2014 kejar L1/L2.",
            "Pola: Cache-Aside (lazy, baca dulu cache, miss baru DB) vs Write-Through (tulis cache+DB bersamaan).",
            "TTL: jadwal sholat 1 jam, profil 5 menit, komunitas 10 menit, cari populer 5 menit \u2014 jangan cache selamanya.",
            "Tiering: hot (sering diakses, TTL pendek), warm (kadang, TTL menengah), cold (jarang, tanpa cache).",
            "Kapan Redis wajib? p99 cache >5ms atau >500 writes/detik atau hot data >80% hit rate (threshold 80%).",
        ],
        table_headers=["Tier", "Latency", "Contoh", "TTL"],
        table_rows=[
            ["L1 memory", "sub-ms", "Feature flag, JWT cache", "60s"],
            ["L2 Redis", "1-5ms", "Jadwal sholat, profil", "5m-1h"],
            ["L3 Postgres", "10-50ms", "Transaksi, ledger", "-"],
        ],
    )
)

# 19 Bab5 ES
slides.append(
    dict(
        title="Bab 5  \u2014  Elasticsearch: Inverted Index + Geospasial",
        footer="Bab 5",
        bullets=[
            "Inverted Index vs LIKE: LIKE scan semua baris; ES pecah kata \u2192 index kata \u2192 cari kata <10ms untuk jutaan dokumen.",
            "Kapan butuh ES? Fase 3 (500+ komunitas, 50k+ anggota) \u2014 MVP cukup pg_trgm; jangan over-engineering di awal.",
            "Geospasial masjid terdekat: geo_distance 5km + sort _geo_distance <10ms (PostGIS ~50ms) \u2014 butuh lat/lng di communities.",
            "Mapping: text analyzer indonesian + keyword + suggest, geo_point lat_lng, date created_at, keyword kelurahan.",
            "Sinkronisasi: Postgres WAL \u2192 Debezium \u2192 Kafka \u2192 ES (real-time) atau batch 5 menit untuk MVP.",
        ],
    )
)

# 20 Bab6 CDC
slides.append(
    dict(
        title="Bab 6  \u2014  CDC: Debezium WAL \u2192 Kafka \u2192 ES/ClickHouse",
        footer="Bab 6",
        bullets=[
            "CDC = tangkap perubahan DB tanpa polling \u2014 Debezium baca WAL logical (INSERT/UPDATE/DELETE) \u2192 kirim ke Kafka.",
            "Kafka sebagai broker: consumer terpisah tulis ke ES (search) & ClickHouse (OLAP) \u2014 buffer & replay.",
            "Batch ETL alternatif untuk non real-time: cron 00:00 ekstrak Postgres \u2192 ClickHouse (laporan bulanan, hemat biaya).",
            "Anti dual-write: tulis HANYA ke Postgres (source of truth); jangan insert ke Postgres lalu ke ES di kode aplikasi \u2014 jika salah satu gagal = inkonsisten.",
            "Prinsip single writer: Postgres WAL adalah kebenaran; ES/ClickHouse adalah derived (bisa rebuild kapan saja).",
        ],
        mono="Postgres WAL --(Debezium pgoutput)--> Kafka(9092) --+--> ES(9200) geo <10ms\n                                           +--> ClickHouse(8123) OLAP\nDual-write DANGER: app->PG OK + app->ES FAIL = data hilang di search",
    )
)

# 21 Bab7 API
slides.append(
    dict(
        title="Bab 7  \u2014  API Delivery: Payload, GZIP, Cursor, Edge, Rate Limit",
        footer="Bab 7",
        bullets=[
            "Payload shaping: ?fields=name,lat,lng \u2014 kirim hanya yang dibutuhkan, hemat 70% bandwidth & parsing.",
            "GZIP 70-80% (semua browser), Brotli +20-30% lebih baik \u2014 aktifkan di Gateway (level 6, threshold 1024).",
            "Cursor vs OFFSET: cursor 20ms stabil, OFFSET 2000ms di halaman dalam \u2014 wajib cursor untuk feed/list.",
            "Cloudflare Edge caching: cache response GET di edge (TTL 1m) \u2014 kurangi origin hit, TTFB <50ms global.",
            "Rate limiting: 100 req/menit (umum), 10 (write), 5 (lapor) \u2014 lindungi backend dari burst & abuse.",
        ],
    )
)

# 22 Bab8 Observability
slides.append(
    dict(
        title="Bab 8  \u2014  Observabilitas: 3 Pilar + pg_stat_statements",
        footer="Bab 8",
        bullets=[
            "pg_stat_statements: lacak query lambat (mean_time, calls) \u2014 SELECT * FROM pg_stat_statements ORDER BY mean_time DESC.",
            "Slow log 100ms: log_min_duration_statement = 100ms \u2192 Loki; identifikasi bottleneck tanpa tebak.",
            "Prometheus (9090) + Grafana (3000): 10 metrik Poster #6, 11 scrape jobs, interval 15s, histogram_quantile p95/p99.",
            "Alert threshold: p95 >500ms 5m, error >1%, DB pool >80%, CPU >80%, cache hit <80% \u2192 WA/Email via Alertmanager.",
            "3 pilar: Metrics (apa yang lambat), Logs (kenapa, via Loki LogQL requestId), Traces (di mana, via Jaeger flame graph).",
        ],
    )
)

# 23 Bab9 Roadmap
slides.append(
    dict(
        title="Bab 9  \u2014  Roadmap 5 Fase: MVP Rp0 \u2192 Auto-scale",
        footer="Bab 9",
        bullets=[
            "MVP (Bln 1-6) Rp0: Postgres + pg_trgm + MatView + PgBouncer \u2192 p50 <50ms, 100 RPS, 5-10 komunitas.",
            "Fase 2 (7-12) Rp0-500rb: + Redis caching + pooling \u2192 500 RPS, 50+ komunitas, Supabase Pro + Upstash Free.",
            "Fase 3 (13-24) Rp1.5-5jt: + ES + CDC + ClickHouse \u2192 5.000 RPS, 500+ komunitas, geospasial & analitik.",
            "Fase 4 (25+) Rp10jt+: + Full observability + HA/DR \u2192 50.000 RPS, 5.000+ komunitas, 99.9%.",
            "Fase 5 (25+) Auto-scale: p50 <10ms, 200.000+ RPS, 50.000+ komunitas, 99.99% \u2014 mulai sederhana, tingkatkan bertahap.",
        ],
        table_headers=["Fase", "Throughput", "Biaya", "p50"],
        table_rows=[
            ["MVP", "100 RPS", "Rp 0", "<50ms"],
            ["Fase 2", "500 RPS", "Rp 0-500rb", "<50ms"],
            ["Fase 3", "5.000 RPS", "Rp 1.5-5jt", "<30ms"],
            ["Fase 4", "50.000 RPS", "Rp 10jt+", "<20ms"],
            ["Fase 5", "200.000 RPS", "Auto-scale", "<10ms"],
        ],
    )
)

# 24 Bab10.1 SLA
slides.append(
    dict(
        title="Bab 10.1  \u2014  SLA 16 Endpoint (p50 / p95 / p99)",
        footer="Bab 10",
        bullets=[
            "16 endpoint dengan SLA per-endpoint \u2014 tidak ada satu angka untuk semua; baca tabel sebelum rilis.",
            "Read tercepat: /jadwal-sholat p50 <20ms (Redis TTL 1 jam); Write: /kas POST p50 <100ms (trigger SHA-256).",
            "Semua endpoint wajib EXPLAIN ANALYZE sebelum rilis; jika p95 >500ms \u2192 incident.",
            "Strategi utama per endpoint: cache, index, MatView, cursor, GIN \u2014 lihat kolom Strategi di tabel.",
        ],
        table_headers=["#", "Endpoint", "p50", "p95", "p99"],
        table_rows=[
            ["1", "GET /komunitas/:id", "<30ms", "<100ms", "<200ms"],
            ["3", "GET /kas", "<50ms", "<200ms", "<500ms"],
            ["5", "GET /pengumuman", "<30ms", "<100ms", "<200ms"],
            ["11", "GET /jadwal-sholat", "<20ms", "<50ms", "<100ms"],
            ["12", "GET /cari", "<50ms", "<200ms", "<500ms"],
            ["...", "16 endpoint total", "<50 read", "<200", "<500"],
        ],
    )
)

# 25 Bab10.2 Throughput
slides.append(
    dict(
        title="Bab 10.2  \u2014  Throughput per Fase + Biaya",
        footer="Bab 10",
        bullets=[
            "Rumus: Throughput = (DAU x req_per_user) / (peak_hours x 3600), burst 5x, target 10x estimasi (safety margin).",
            "MVP 100 RPS (500 DAU) \u2192 Fase 2 500 \u2192 Fase 3 5.000 \u2192 Fase 4 50.000 \u2192 Fase 5 200.000+ RPS.",
            "Biaya: Rp0 (Free) \u2192 Rp0-500rb (Pro) \u2192 Rp1.5-5jt (ES+CH) \u2192 Rp10jt+ (Enterprise) \u2014 Rp0-friendly via Podman lokal.",
            "Scaling: MVP-Fase2 vertical (tambah RAM), Fase3 horizontal (+server+LB), Fase4-5 auto-scale.",
            "Koneksi DB: 10-20 \u2192 50-100 \u2192 200-500 \u2192 500-1000 \u2192 2000+ (PgBouncer pool 25 jaga).",
        ],
        table_headers=["Fase", "Komunitas", "DAU", "Target RPS", "Biaya"],
        table_rows=[
            ["MVP", "5-10", "500-1k", "100", "Rp 0"],
            ["Fase 2", "50+", "5k-10k", "500", "Rp 0-500rb"],
            ["Fase 3", "500+", "50k-100k", "5.000", "Rp 1.5-5jt"],
            ["Fase 4", "5.000+", "500k-1jt", "50.000", "Rp 10jt+"],
            ["Fase 5", "50.000+", "5jt+", "200.000+", "Auto"],
        ],
    )
)

# 26 Bab10.4 Checklist
slides.append(
    dict(
        title="Bab 10.4  \u2014  Checklist 10 Definition of Done",
        footer="Bab 10",
        bullets=[
            "Wajib diverifikasi sebelum setiap rilis \u2014 penanggung jawab rilis harus tanda tangan checklist.",
            "Setiap item tidak terpenuhi = risiko degradasi performa; jangan deploy jika ada yang merah.",
            "Alat: psql EXPLAIN, pg_stat_statements, curl GZIP, k6, Grafana p95, code review N+1.",
            "Otomatisasi: CI jalankan EXPLAIN + lint + k6 smoke sebelum merge.",
        ],
        table_headers=["#", "Checklist", "Verifikasi"],
        table_rows=[
            ["1", "EXPLAIN ANALYZE semua query baru", "Execution Time OK"],
            ["2", "No Seq Scan >1000 rows", "EXPLAIN tanpa Seq Scan"],
            ["3", "Semua FK punya index", "pg_constraint vs pg_index"],
            ["4", "pg_trgm aktif jika ada search", "pg_extension"],
            ["5", "MatView refresh terjadwal", "cron 5m / CONCURRENTLY"],
            ["6", "Autovacuum normal", "n_dead_tup <1000"],
            ["7", "p95 <500ms semua endpoint", "histogram_quantile"],
            ["8", "No N+1 query", "JOIN / IN, OTEL trace"],
            ["9", "GZIP/Brotli aktif", "Content-Encoding: gzip"],
            ["10", "Rate limiting aktif", "429 setelah burst"],
        ],
    )
)

# 27 Demo 01
slides.append(
    dict(
        title="Demo 01  \u2014  console.log Anti-Pattern",
        footer="Demo 01",
        bullets=[
            "Bahaya: console.log tanpa level, password bocor di log, tanpa requestId, Seq Scan tersembunyi.",
            "Contoh: console.log('cari q='+q+' ip='+req.ip) \u2014 password ikut ke log, tidak bisa filter level, tidak ada korelasi.",
            "Dampak: log tidak terstruktur \u2192 Loki tidak bisa query; Seq Scan 2000ms tidak terlihat sampai produksi meledak.",
            "Solusi: ganti Pino JSON (Demo 02) + requestId + redact [Redacted] + EXPLAIN ANALYZE.",
            "Coba: jalankan umkm-service/src/index.ts (branch 01) \u2192 lihat log berantakan & query LIKE '%ayam%' Seq Scan.",
        ],
        mono="console.log('query: SELECT * FROM umkm WHERE name LIKE \\'%'+q+'%\\''); // Seq Scan!\nconsole.log('password='+req.body.password); // BOCOR!\n// Tidak ada level, tidak ada requestId, tidak bisa filter di Loki",
    )
)

# 28 Demo 02
slides.append(
    dict(
        title="Demo 02  \u2014  Pino JSON Proper",
        footer="Demo 02",
        bullets=[
            "Pino: level (debug/info/warn/error), redact [Redacted] untuk password/token, requestId, latency_ms, hostname.",
            "Contoh log JSON: {level:'info', service:'umkm-service', requestId:'550e...', latency_ms:23, q:'ayam'} \u2014 terstruktur, bisa di-query Loki.",
            "Redact: paths ['password','token','req.headers.authorization'] \u2192 '[Redacted]' \u2014 aman, tidak bocor.",
            "Transport: pino-pretty di dev (colorize), JSON murni di production \u2192 Alloy \u2192 Loki.",
            "Coba: index-proper.ts \u2192 curl /api/cari?q=ayam \u2192 lihat log JSON + X-Request-Id header.",
        ],
        mono='{"level":"info","service":"umkm-service","requestId":"550e8400...","latency_ms":23,"q":"ayam"}\n{"level":"info","password":"[Redacted]","token":"[Redacted]"} // aman',
    )
)

# 29 Demo 03
slides.append(
    dict(
        title="Demo 03  \u2014  Scale: GIN, Cursor, MatView, RLS",
        footer="Demo 03",
        bullets=[
            "GIN: LIKE '%ayam%' 2000ms \u2192 GIN 10-50ms (40-200x) \u2014 pg_trgm + similarity >0.3.",
            "Cursor: OFFSET 10000 2000ms \u2192 cursor keyset 20ms (100x) \u2014 WHERE (created_at,id) > cursor LIMIT 20.",
            "MatView: SUM(amount) 500ms \u2192 mv_kas_total 5-30ms (16-100x) \u2014 REFRESH CONCURRENTLY tiap 5 menit.",
            "RLS: overhead <0.1ms \u2014 policy community_isolation + index FK, isolasi per komunitas tanpa lambat.",
            "Coba: scripts/explain-demo.sql \u2192 EXPLAIN (ANALYZE, BUFFERS) lihat Index Scan vs Seq Scan.",
        ],
        table_headers=["Teknik", "Sebelum", "Sesudah", "Speedup"],
        table_rows=[
            ["GIN pg_trgm", "2000ms", "10-50ms", "40-200x"],
            ["Cursor", "2000ms", "20ms", "100x"],
            ["MatView", "500ms", "5-30ms", "16-100x"],
            ["RLS", "-", "<0.1ms", "-"],
        ],
    )
)

# 30 Demo 04
slides.append(
    dict(
        title="Demo 04  \u2014  Observability: Grafana, Loki, Jaeger",
        footer="Demo 04",
        bullets=[
            "Grafana 10 metrik (Poster #6) + per-endpoint 16 SLA + proteksi (Circuit/Bulkhead/Backpressure) \u2014 dashboard performa-gr.",
            'Loki LogQL: {service=~"umkm|kas"} | json | requestId="550e..." \u2014 korelasi log lintas service via requestId.',
            "Jaeger flame graph: trace 320ms \u2192 span pg.query 40ms, redis 5ms, http 280ms \u2014 langsung tahu bottleneck.",
            "Alert: P95High >500ms 5m, ErrorRate >1%, DBPool >80%, CacheHit <80% \u2192 WA/Email via Alertmanager.",
            "Coba: podman-compose --profile observability up -d \u2192 open :3000 (admin/admin), :9090, :16686, :3100.",
        ],
    )
)

# 31 Demo 05
slides.append(
    dict(
        title="Demo 05  \u2014  CDC Streaming: Debezium \u2192 Kafka \u2192 ES <10ms",
        footer="Demo 05",
        bullets=[
            "Debezium baca WAL logical (slot debezium, publication dbz_publication) \u2192 Kafka topic gotongroyong.* \u2192 consumer tulis ES/ClickHouse.",
            "ES geo_distance 5km <10ms untuk 256 masjid (vs PostGIS ~50ms) \u2014 demo: GET /api/masjid-terdekat?lat=-6.25&lng=106.75&radius=5km.",
            "Full-text multi_match + highlight + fuzziness AUTO \u2192 GET /api/cari-es?q=ayam, fallback pg_trgm jika ES down.",
            "ClickHouse: ledger_analytics MergeTree ORDER BY (community_id, created_at) + mv_ledger_daily SummingMergeTree.",
            "Anti dual-write: app hanya tulis Postgres; CDC yang sebar \u2014 lag_ms dimonitor, alert jika >1000ms.",
        ],
        mono='curl "http://localhost:3003/api/masjid-terdekat?lat=-6.25&lng=106.75&radius=5km" | jq .meta.took_ms\n# expect <10  (ES geo_distance)\nbash scripts/es-demo.sh  # create index, bulk 6k, verify',
    )
)

# 32 Appendix A
slides.append(
    dict(
        title="Appendix A  \u2014  5 Proteksi: Token, Circuit, Bulkhead, Backpressure, Degradation",
        footer="Appendix A",
        bullets=[
            "Token/Leaky Bucket: rate limiting 100/10/5 req/menit \u2014 tolak burst, kembalikan 429, simpan di Redis.",
            "Circuit Breaker: jika downstream gagal 50% 10 detik \u2192 open 30 detik, fallback cache/error cepat, bukan hang.",
            "Bulkhead: isolasi pool per service (max 25) \u2014 satu service lambat tidak habiskan semua koneksi.",
            "Backpressure: queue penuh \u2192 tolak dengan 503 + Retry-After, jangan OOM.",
            "Graceful Degradation: jika ES down \u2192 fallback pg_trgm; jika Redis down \u2192 langsung DB (lebih lambat tapi tetap hidup).",
        ],
    )
)

# 33 Appendix B
slides.append(
    dict(
        title="Appendix B  \u2014  Threshold Scaling: 1TB / 10M / 1000",
        footer="Appendix B",
        bullets=[
            "Threshold Shopee \u2192 GR konservatif: Data >500GB (vs 1TB), Single table >5M rows (vs 10M), Write QPS >500 (vs 1000) \u2192 evaluasi sharding/replica.",
            "Vertical vs Horizontal: <1000 RPS vertical (tambah RAM/CPU) cukup; >1000 RPS horizontal (+server+LB); >50k auto-scale.",
            "Replica/Sharding: read replica untuk baca, sharding by community_id untuk write-heavy, MatView untuk agregasi.",
            "HA/DR: replica Postgres, backup WAL, failover otomatis, RTO <5 menit, RPO <1 menit (WAL archiving).",
            "Monitor 80% threshold, alert sebelum 100% \u2014 jangan tunggu penuh baru panik.",
        ],
        table_headers=["Metrik", "Threshold GR", "Tindakan"],
        table_rows=[
            ["Data size", ">500 GB", "Sharding / replica"],
            ["Single table", ">5M rows", "Partitioning, archiving"],
            ["Write QPS", ">500/detik", "Kafka queue, batch"],
            ["P99 cache", ">5ms", "Wajib Redis"],
        ],
    )
)

# 34 Appendix C
slides.append(
    dict(
        title="Appendix C  \u2014  Security 10 Layer",
        footer="Appendix C",
        bullets=[
            "10 layer: TLS 1.3 (transit), AES-256 (at-rest), RLS (row isolation), hash chain SHA-256 (ledger), Vault (secret), RBAC, rate limit, audit log, backup encrypt, WAF.",
            "TLS 1.3 wajib untuk semua HTTP; AES-256 untuk storage S3/R2; RLS <0.1ms untuk multi-tenant.",
            "Hash chain: hash_self = SHA256(amount|description|recipient|actor|hash_prev), genesis 0x00..00, verifikasi via GET /ledger/verify.",
            "Vault: jangan hardcode secret di env file \u2014 ambil dari Vault/KMS, rotasi otomatis.",
            "Audit: tabel audit_log (actor, action, timestamp, ip, old/new JSON) + trigger \u2014 siapa ubah apa, kapan.",
        ],
    )
)

# 35 Penutup
slides.append(
    dict(
        title="Penutup  \u2014  Kecepatan = Amanah",
        footer="Penutup",
        cover=True,
        cover_title="Kecepatan = Amanah",
        cover_sub="Mulai sederhana, tingkatkan bertahap \u2014 setiap ms = kepercayaan",
        cover_bullets=[
            "MVP Rp0 sudah p50 <50ms dengan Postgres+Redis+pg_trgm+MatView",
            "Tambah ES+CDC+ClickHouse saat 500+ komunitas (Fase 3) \u2014 jangan over-engineering di awal",
            "Ukur, jangan tebak: EXPLAIN ANALYZE, pg_stat_statements, Prometheus, Loki, Jaeger",
            "TIGA INSAN: Muttaqin (percaya), Shalih (berkarya), Nafi' (bermanfaat) \u2014 performa adalah wujudnya",
        ],
    )
)


# ──────────────────────────────────────────────
# RENDER
# ──────────────────────────────────────────────
def render_cover(slide, data, idx):
    add_bg(slide)
    # navy full cover bg
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    # orange accent bottom
    acc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        prs.slide_height - Inches(0.18),
        prs.slide_width,
        Inches(0.18),
    )
    acc.fill.solid()
    acc.fill.fore_color.rgb = ORANGE
    acc.line.fill.background()
    # logo text top left
    tx_logo = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(4), Inches(0.3)
    )
    tf = tx_logo.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "GOTONG ROYONG  \u2022  OS Kehidupan Komunitas"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0xFD, 0xBA, 0x74)
    p.font.name = "Calibri"
    p.font.bold = True
    # title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(7.5), Inches(1.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data["cover_title"]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    # subtitle
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(7.5), Inches(0.5))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = data["cover_sub"]
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p2.font.name = "Calibri"
    # bullets card on right
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.3), Inches(4.5), Inches(4.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()
    card.shadow.inherit_shadow = False
    tx3 = slide.shapes.add_textbox(Inches(8.5), Inches(1.55), Inches(4.0), Inches(3.7))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    for i, b in enumerate(data["cover_bullets"]):
        p = tf3.add_paragraph() if i > 0 else tf3.paragraphs[0]
        p.text = "\u2022  " + b
        p.font.size = Pt(9)
        p.font.color.rgb = NAVY
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    # footer
    txf = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.55), Inches(7), Inches(0.3)
    )
    tff = txf.text_frame
    tff.word_wrap = True
    p = tff.paragraphs[0]
    p.text = "60 menit  \u00b7  10 Bab  \u00b7  16 Endpoint  \u00b7  Checklist 10 DoD  \u00b7  Rp0 \u2192 Auto-scale"
    p.font.size = Pt(7)
    p.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    p.font.name = "Calibri"
    # slide number
    txn = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.2),
        prs.slide_height - Inches(0.55),
        Inches(1),
        Inches(0.3),
    )
    tfn = txn.text_frame
    tfn.word_wrap = True
    p = tfn.paragraphs[0]
    p.text = f"{idx} / 35"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(7)
    p.font.color.rgb = ORANGE
    p.font.name = "Calibri"
    p.font.bold = True


NOTES = {
    1: "Modul Performa Backend GR Demo\nNarasi: 4 branch + 5 CDC, TIGA INSAN Muttaqin, p50<50ms MVP Rp0\nQ&A: Kenapa Rp0? Podman lokal tanpa cloud",
    2: "Daftar Isi -- 10 Bab + 35 Slides Map\nNarasi: Peta 10 bab, Poster 1-7 di part1 bangun bahasa sebelum teknik\nQ&A: Harus hafal 35? Tidak, cukup peta",
    3: "Bab 1.1 -- Kecepatan = Kepercayaan (Muttaqin)\nNarasi: Lambat=khianat amanah, SHA-256 <1 detik, ukur p50/p95/p99\nQ&A: Hubungan Muttaqin-ms? Amanah = sistem cepat transparan",
    4: "Bab 1.2 -- Konteks Indonesia 3G 2-3GB RAM\nNarasi: RTT 500-1000ms, backend <200ms agar total <3 detik, GZIP 70%\nQ&A: Kenapa <200ms? Sisa budget untuk 3G",
    5: "Bab 1.3 -- Target SLA p50/p95/p99\nNarasi: p50<50ms read, p95<200ms, p99<500ms, 99.5%->99.99%, prom-client histogram\nQ&A: Kenapa read/write beda? Write ACID+trigger lebih lambat",
    6: "Poster #1 -- Apa itu 200ms\nNarasi: 1 detik=1000ms, 200ms=0.2 detik kedipan mata, User->Server->DB->Response 30-50ms\nPerintah Live: Tidak ada demo, fokus teori\nQ&A: Kenapa 200ms? Backend <200ms agar total dengan 3G <3 detik",
    7: "Poster #2 -- Semakin kecil ms semakin cepat\nNarasi: 1ms sangat cepat (Redis HIT) vs 2000ms sangat lambat (Seq Scan) 400-2000x\nQ&A: 1ms mungkin? Ya Redis HIT & RLS <0.1ms",
    8: "Poster #3 -- Apa itu P99\nNarasi: 100 request urut, P99=request ke-99, avg 109ms bohong jika 1 request 10000ms\nQ&A: Kenapa bukan rata-rata? Tertarik outlier, P99 jujur",
    9: "Poster #4 -- P50/P95/P99/P99.9\nNarasi: P50 median 50%, P95 5% tail, P99 1% tail, P99.9 0.1% untuk 200k RPS\nQ&A: Kapan P99.9 penting? Fase5 200k RPS, 0.1%=200 req/detik lambat",
    10: "Poster #5 -- Skala rasa\nNarasi: 0-50 sangat cepat, 50-100 sangat baik, 100-200 baik, 200-500 mulai terasa, >1000 sangat lambat\nQ&A: Skala server atau total? Server-side, total +RTT 500ms",
    11: "Poster #6.1 -- 10 Metrik Wajib\nNarasi: ResponseTime, P95/P99, RPS, Latency, Error<0.1%, Availability, CPU<70%, Memory<70%, DB<50ms, Cache>80%\nQ&A: Harus 10 sekaligus? Mulai 3: ResponseTime, Error, CacheHit",
    12: "Poster #7 -- Glossary Super Sederhana\nNarasi: ms, P99, RPS, Latency, Throughput, Error, Availability, CPU, Cache Hit -- bahasa sama\nQ&A: Bedanya latency vs throughput? Waktu satu request vs kapasitas per detik",
    13: "Bab 2 -- Data Flow Flutter->Gateway->7 Fondasi->6 DB\nNarasi: Flutter->Gateway JWT<5ms->7 Fondasi->6 DB, hot path Redis HIT <5ms\nPerintah Live: Tidak ada demo, fokus arsitektur\nQ&A: Kenapa 7 fondasi? Dipakai semua modul, hindari duplikasi",
    14: "Bab 3.1 -- Index B-Tree 50.000x\nNarasi: Seq Scan 1M langkah 2000ms vs B-Tree 20 langkah 10ms O(log n), FK wajib index CONCURRENTLY\nPerintah Live: psql EXPLAIN ANALYZE SELECT * FROM umkm WHERE kelurahan='Bintaro' -> Index Scan 10ms\nQ&A: Kapan index merugikan? >50% full scan atau >5-7 index per tabel",
    15: "Bab 3.2 -- pg_trgm GIN 10-50ms vs LIKE 2000ms\nNarasi: LIKE '%ayam%' Seq Scan 2000ms, pg_trgm trigram GIN 10-50ms 40-200x, WHERE name % 'ayam'\nPerintah Live: curl /api/cari?q=ayam 10-50ms vs ?q=ayam&mode=like 2000ms\nQ&A: Kenapa GIN bukan B-Tree? B-Tree hanya =/prefix, GIN untuk LIKE & %",
    16: "Bab 3.3-3.4 -- MatView + Cursor vs OFFSET\nNarasi: MatView SUM 500ms->5-30ms 16-100x refresh CONCURRENTLY 5m, OFFSET 10000 2000ms->cursor 20ms 100x\nPerintah Live: curl /api/kas?mode=before 500ms vs 5-30ms; offset=10000 2000ms vs cursor 20ms\nQ&A: Kapan cursor vs OFFSET? Cursor infinite scroll, OFFSET hanya halaman kecil <100",
    17: "Bab 3.5-3.8 -- PgBouncer, EXPLAIN, RLS, VACUUM\nNarasi: PgBouncer pool25 hemat 95% RAM, EXPLAIN 6 langkah, RLS <0.1ms, VACUUM n_dead_tup<1000\nPerintah Live: psql :6432 SHOW POOLS; EXPLAIN ANALYZE ledger; SELECT n_dead_tup\nQ&A: Kenapa pool 25? (CPU*2)+spindle, 4 core=10-15 + buffer 100 RPS",
    18: "Bab 4 -- Caching L1/L2/L3 + TTL + Tiering\nNarasi: L1 sub-ms, L2 Redis 1-5ms, L3 PG 10-50ms, Cache-Aside, TTL sholat 1j profil 5m\nPerintah Live: curl -i /api/komunitas/xxx grep X-Cache MISS lalu HIT; redis-cli INFO stats hit>80%\nQ&A: Kenapa Cache-Aside? Resilience, Redis down tetap jalan via DB",
    19: "Bab 5 -- Elasticsearch Inverted Index + Geospasial\nNarasi: Inverted <10ms vs LIKE 2000ms, geo_distance 5km <10ms vs PostGIS 50ms, MVP pg_trgm cukup\nPerintah Live: curl :9200/umkm/_search?q=ayam took<10ms; /masjid-terdekat?lat=-6.25&lng=106.75 <10ms\nQ&A: Kenapa tidak ES dari MVP? pg_trgm gratis cukup puluhan ribu, ES Fase3 500+ komunitas",
    20: "Bab 6 -- CDC Debezium WAL->Kafka->ES/ClickHouse\nNarasi: WAL logical -> Kafka -> ES/CH, anti dual-write single writer PG, lag_ms<1000ms\nPerintah Live: kafka-console-consumer --topic gotongroyong.public.umkm | jq .payload; INSERT umkm cek ES\nQ&A: Kenapa tidak dual-write? PG OK+ES FAIL=inkonsisten, tidak ada transaksi lintas DB",
    21: "Bab 7 -- API Delivery Payload GZIP Cursor Edge RateLimit\nNarasi: ?fields= hemat 70%, GZIP 70-80% Brotli +20%, cursor 20ms vs OFFSET 2000ms, Edge TTL1m, rate 100/10/5\nPerintah Live: curl -H Accept-Encoding:gzip -v /api/umkm?fields=name,lat,lng grep gzip 20KB vs 100KB\nQ&A: Kenapa sparse fieldsets? 15KB vs 100KB hemat kuota 3G & parsing 2GB RAM",
    22: "Bab 8 -- Observabilitas 3 Pilar + pg_stat_statements\nNarasi: pg_stat_statements top total_exec_time, slow log 100ms->Loki, Prometheus 9090 11 jobs, alert p95>500ms 5m\nPerintah Live: psql SELECT * FROM pg_stat_statements ORDER BY total_exec_time DESC; curl :9090/api/v1/query?query=up\nQ&A: Kenapa slow log 100ms? Sweet spot, 10ms banjir, 1000ms lewatkan pola",
    23: "Bab 9 -- Roadmap 5 Fase MVP Rp0->Auto-scale\nNarasi: MVP 100 RPS Rp0 PG+pg_trgm+MatView, F2 500 RPS Redis, F3 5k ES+CDC Rp1.5-5jt, F4 50k Rp10jt+, F5 200k 99.99%\nPerintah Live: Tidak ada demo, fokus roadmap visual\nQ&A: Kenapa tidak langsung Fase5? Biaya Rp10jt+ & trafik belum ada, mulai sederhana bertahap",
    24: "Bab 10.1 -- SLA 16 Endpoint p50/p95/p99\nNarasi: /jadwal-sholat p50<20ms Redis1j tercepat, /kas p50<50ms MatView, /cari p50<50ms GIN/ES, ukur X-Response-Time\nPerintah Live: k6 run load/k6-script.js; curl -i /api/komunitas/xxx grep X-Response-Time\nQ&A: Kenapa server-side? Tanpa RTT 3G 500-1000ms agar kontrak stabil",
    25: "Bab 10.2 -- Throughput per Fase + Biaya\nNarasi: Rumus (DAU*req)/(peak*3600)*10, MVP 100 RPS 500 DAU Rp0 -> F5 200k 5jt DAU auto-scale\nPerintah Live: TARGET=:3003 bun --cwd load run load.ts --vus 100; k6 run --vus 100 --duration 30s\nQ&A: Kenapa biaya naik Rp0->Rp10jt+? Infra dedicated & multi-region, revenue marketplace/BMT cover",
    26: "Bab 10.4 -- Checklist 10 Definition of Done\nNarasi: EXPLAIN, no SeqScan>1000, FK index, pg_trgm, MatView 5m, autovacuum <1000, p95<500ms, no N+1, GZIP, rate limit 429 -- 10/10 hijau baru rilis\nPerintah Live: psql -f scripts/explain-demo.sql; SELECT n_dead_tup; curl -H Accept-Encoding:gzip -I grep gzip; burst 105 cek 429\nQ&A: Kenapa tiap sprint? Cegah utang performa, jangan cek sekali di akhir",
    27: 'Demo 01 -- console.log Anti-Pattern\nNarasi: console.log tanpa level, password bocor, tanpa requestId, LIKE Seq Scan 2000ms, avg bohong\nPerintah Live: node order-service/src/index.ts & curl POST /checkout d=\'{"password":"rahasia123"}\' lihat bocor\nQ&A: Bahaya console.log? Tidak terstruktur, Loki tidak bisa query, Seq Scan tersembunyi',
    28: 'Demo 02 -- Pino JSON Proper\nNarasi: Pino level+redact [Redacted]+requestId+latency_ms, JSON terstruktur Loki, pino-pretty dev\nPerintah Live: LOG_LEVEL=debug bun run src/index-proper.ts & curl -H x-request-id:550e... lihat {"password":"[Redacted]"}\nQ&A: Kenapa redact? UU PDP denda Rp2M, Pino redact otomatis aman',
    29: "Demo 03 -- Scale GIN Cursor MatView RLS\nNarasi: GIN 2000->10ms 200x, Cursor 2000->20ms 100x, MatView 500->30ms 16x, RLS <0.1ms, Buffers 45 vs 1200\nPerintah Live: psql EXPLAIN ANALYZE name % 'ayam' 10ms vs LIKE 2000ms; curl offset=10000 2000ms vs cursor 20ms\nQ&A: Kenapa cursor 100x? Langsung lompat via index (created_at,id)>cursor tanpa baca-buang",
    30: "Demo 04 -- Observability Grafana Loki Jaeger\nNarasi: Grafana 10 metrik+16 SLA+proteksi :3000, Loki LogQL requestId :3100, Jaeger flame 320ms pg 40ms redis 5ms :16686, alert P95High>500ms\nPerintah Live: podman-compose --profile observability up -d; open :3000 admin/admin, :9090, :16686, :3100\nQ&A: Bedanya logs vs traces? Logs kenapa error per service, traces di mana bottleneck flame graph",
    31: "Demo 05 -- CDC Debezium->Kafka->ES <10ms\nNarasi: WAL slot debezium publication dbz_publication -> Kafka gotongroyong.* -> ES geo 5km <10ms vs PostGIS 50ms, fallback pg_trgm\nPerintah Live: curl /api/masjid-terdekat?lat=-6.25&lng=106.75&radius=5km | jq .meta.took_ms expect <10; bash scripts/es-demo.sh\nQ&A: Lag CDC? lag_ms=now-ts_ms alert >1000ms, offset commit manual idempotent",
    32: "Appendix A -- 5 Proteksi Token Circuit Bulkhead Backpressure Degradation\nNarasi: Rate 100/10/5 429, Circuit open 30s jika 50% error 10s, Bulkhead pool 25 isolasi, Backpressure 503 Retry-After, Degradation fallback pg_trgm\nPerintah Live: burst 105 cek 429; trigger 50% error cek 503 circuit open; curl /bulkhead/stats\nQ&A: Kenapa 5 proteksi? Cegah domino failure berbeda, satu saja tidak cukup",
    33: "Appendix B -- Threshold Scaling 1TB/10M/1000\nNarasi: GR konservatif 500GB/5M/500 QPS alert 80% vs Shopee 1TB/10M/1000, vertical <1000 RPS horizontal >1000 auto-scale >50k\nPerintah Live: psql SELECT pg_database_size/1GB; cek size <500GB jika > alert scaling\nQ&A: Kapan sharding? Data>500GB atau table>5M atau write>500/detik evaluasi sharding/replica",
    34: "Appendix C -- Security 10 Layer\nNarasi: TLS1.3, AES-256, RLS <0.1ms, SHA-256 hash chain, Vault, RBAC, rate limit, audit log, backup encrypt, WAF\nPerintah Live: Tidak ada demo, fokus teori\nQ&A: Hash chain? hash_self=SHA256(amount|desc|recipient|actor|hash_prev) trigger BEFORE INSERT",
    35: "Penutup -- Kecepatan = Amanah\nNarasi: MVP Rp0 p50<50ms PG+Redis+pg_trgm+MatView, tambah ES+CDC Fase3 500+ komunitas, ukur EXPLAIN+pg_stat+Prometheus TIGA INSAN Muttaqin Shalih Nafi'\nPerintah Live: Tidak ada demo, Q&A 10 menit\nQ&A: Next step? Clone repo, Podman compose, B-Tree FK, pg_trgm, MatView, PgBouncer 25, ukur p95 tiap sprint",
}

for idx, data in enumerate(slides, start=1):
    slide = prs.slides.add_slide(BLANK)
    if data.get("cover"):
        render_cover(slide, data, idx)
        try:
            notes_slide = slide.notes_slide
            notes_slide.placeholders[1].text = NOTES.get(idx, data.get("title", ""))
        except:
            pass
        continue
    add_bg(slide)
    # header
    bab_label = data.get("footer", "")
    add_header(slide, data["title"], f"{bab_label}  \u00b7  {idx}/35")
    # bullets
    bullets = data.get("bullets", [])
    has_table = "table_headers" in data
    has_mono = "mono" in data
    # layout: if table or mono, bullets left 7.2, right card 5.5
    if has_table or has_mono:
        add_bullets(
            slide,
            bullets,
            left=Inches(0.4),
            top=Inches(1.25),
            width=Inches(6.9),
            height=Inches(4.0),
            font_size=Pt(9.5),
        )
        if has_table:
            add_table_slide(
                slide,
                data["table_headers"],
                data["table_rows"],
                Inches(7.6),
                Inches(1.25),
                Inches(5.3),
                Inches(4.2),
            )
        if has_mono:
            add_mono_box(
                slide,
                data["mono"],
                Inches(7.6),
                Inches(1.25) + (Inches(4.4) if has_table else Inches(0)),
                Inches(5.3),
                Inches(1.6) if has_table else Inches(4.2),
            )
            # if both table and mono, adjust mono position below table
            if has_table:
                # remove previous mono and re-add below table
                # we already added, but need to ensure not overlapping — keep as is, table 4.2 + mono 1.6 = 5.8 fits
                pass
    else:
        # no table/mono — check if many bullets, use wider
        if len(bullets) > 5:
            add_bullets(
                slide,
                bullets,
                left=Inches(0.4),
                top=Inches(1.25),
                width=Inches(12.5),
                height=Inches(5.6),
                font_size=Pt(10),
            )
        else:
            add_bullets(
                slide,
                bullets,
                left=Inches(0.4),
                top=Inches(1.25),
                width=Inches(7.0),
                height=Inches(5.6),
                font_size=Pt(10.5),
            )
            # add visual card on right for some slides
            if idx in [13, 16, 17, 18, 19, 20, 21, 22]:
                icons = {
                    13: "\U0001f5fa Data Flow",
                    16: "\u26a1 MatView+Cursor",
                    17: "\U0001f512 RLS+VACUUM",
                    18: "\U0001f4be Caching",
                    19: "\U0001f50d Elasticsearch",
                    20: "\U0001f504 CDC",
                    21: "\U0001f680 API",
                    22: "\U0001f50e Observability",
                }
                # add decorative card
                card_text = {
                    13: "Flutter \u2192 Gateway \u2192 7 Fondasi \u2192 6 DB\nHot path: Redis HIT <5ms",
                    16: "MatView 500\u219230ms\nCursor 2000\u219220ms",
                    17: "PgBouncer 25 pool\nRLS <0.1ms\nVACUUM rutin",
                    18: "L1 sub-ms\nL2 1-5ms\nL3 10-50ms",
                    19: "Inverted Index\n500+ komunitas\n\u2192 ES",
                    20: "Single writer:\nPG WAL \u2192 Kafka \u2192 ES/CH",
                    21: "GZIP 70-80%\nCursor > OFFSET\nEdge cache",
                    22: "Metrics/Logs/Traces\np95>500ms \u2192 alert",
                }
                add_card(
                    slide,
                    Inches(7.7),
                    Inches(1.4),
                    Inches(5.1),
                    Inches(2.2),
                    icons.get(idx, ""),
                    card_text.get(idx, ""),
                    "",
                )
    add_footer(slide, data.get("footer", ""), idx)
    try:
        notes_slide = slide.notes_slide
        notes_slide.placeholders[1].text = NOTES.get(idx, data.get("title", ""))
    except:
        pass

out = "/home/ngome/GotongRoyong/backend-performa-demo/presentasi/Modul_Performa_Backend_GR_Demo.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides)} slides")
