#!/usr/bin/env python3
# presentasi/generate_light.py — Light 10 slides backup 20 menit
# Filter dari generate_pptx.py slides: [0 Cover, 1 Daftar Isi, 5 Poster #1, 7 Poster #3 P99, 13 B-Tree, 14 pg_trgm, 17 Caching L1/L2/L3, 19 CDC, 23 SLA 16 endpoint, 34 Penutup]
# Spec tulis 21 CDC / 27 SLA (1-based slide num), koreksi 0-based 19/23 agar semantik tepat
import pathlib
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
NAVY_DARK = RGBColor(0x14, 0x26, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)

prs = pptx.Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()


def add_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.95)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    acc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.95), prs.slide_width, Inches(0.07)
    )
    acc.fill.solid()
    acc.fill.fore_color.rgb = ORANGE
    acc.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.5), Inches(0.65))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    if subtitle:
        tx2 = slide.shapes.add_textbox(
            Inches(10.0), Inches(0.28), Inches(3.0), Inches(0.4)
        )
        tf2 = tx2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.alignment = PP_ALIGN.RIGHT
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
        p2.font.name = "Calibri"


def add_footer(slide, bab, idx, total=10):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        prs.slide_height - Inches(0.32),
        prs.slide_width,
        Inches(0.32),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY_DARK
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(
        Inches(0.3), prs.slide_height - Inches(0.28), Inches(9.5), Inches(0.22)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Gotong Royong — OS Kehidupan Komunitas  |  60 menit  |  {bab}  |  TIGA INSAN: Muttaqin • Shalih • Nafi'"
    p.font.size = Pt(7)
    p.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p.font.name = "Calibri"
    tx2 = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.2),
        prs.slide_height - Inches(0.28),
        Inches(1.0),
        Inches(0.22),
    )
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f"{idx} / {total}"
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(7)
    p2.font.color.rgb = RGBColor(0xFD, 0xBA, 0x74)
    p2.font.name = "Calibri"
    p2.font.bold = True


def add_bullets(
    slide,
    bullets,
    left=Inches(0.5),
    top=Inches(1.35),
    width=Inches(7.2),
    height=Inches(5.4),
    font_size=Pt(11),
):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "•  " + b
        p.font.size = font_size
        p.font.name = "Calibri"
        p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        p.space_after = Pt(6)
        p.space_before = Pt(2)
    return tx


def add_mono_box(slide, text, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    shape.line.fill.background()
    tx = slide.shapes.add_textbox(
        left + Inches(0.12),
        top + Inches(0.08),
        width - Inches(0.24),
        height - Inches(0.16),
    )
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x86, 0xEF, 0xAC)


def add_table_slide(slide, headers, rows, left, top, width, height):
    tbl_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, width, height
    )
    tbl = tbl_shape.table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(7)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(7)
                    run.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
                    run.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(0xF1, 0xF5, 0xF9) if i % 2 == 0 else WHITE
            )
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl_shape


def render_cover(slide, data, idx, total=10):
    add_bg(slide)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    acc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        prs.slide_height - Inches(0.18),
        prs.slide_width,
        Inches(0.18),
    )
    acc.fill.solid()
    acc.fill.fore_color.rgb = ORANGE
    acc.line.fill.background()
    tx_logo = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(4), Inches(0.3)
    )
    tf = tx_logo.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "GOTONG ROYONG  •  OS Kehidupan Komunitas"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0xFD, 0xBA, 0x74)
    p.font.name = "Calibri"
    p.font.bold = True
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(7.5), Inches(1.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data["cover_title"]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(7.5), Inches(0.5))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = data["cover_sub"]
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p2.font.name = "Calibri"
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.3), Inches(4.5), Inches(4.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()
    tx3 = slide.shapes.add_textbox(Inches(8.5), Inches(1.55), Inches(4.0), Inches(3.7))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    for i, b in enumerate(data["cover_bullets"]):
        p = tf3.add_paragraph() if i > 0 else tf3.paragraphs[0]
        p.text = "•  " + b
        p.font.size = Pt(9)
        p.font.color.rgb = NAVY
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    txf = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.55), Inches(7), Inches(0.3)
    )
    tff = txf.text_frame
    tff.word_wrap = True
    p = tff.paragraphs[0]
    p.text = (
        "60 menit  ·  10 Bab  ·  16 Endpoint  ·  Checklist 10 DoD  ·  Rp0 → Auto-scale"
    )
    p.font.size = Pt(7)
    p.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    p.font.name = "Calibri"
    txn = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.2),
        prs.slide_height - Inches(0.55),
        Inches(1),
        Inches(0.3),
    )
    tfn = txn.text_frame
    tfn.word_wrap = True
    p = tfn.paragraphs[0]
    p.text = f"{idx} / {total}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(7)
    p.font.color.rgb = ORANGE
    p.font.name = "Calibri"
    p.font.bold = True


# Coba import slides dari generate_pptx.py jika modular, fallback manual
LIGHT_IDX = [0, 1, 5, 7, 13, 14, 17, 19, 23, 34]
try:
    import importlib.util, pathlib

    spec = importlib.util.spec_from_file_location(
        "gen_full", str(pathlib.Path(__file__).with_name("generate_pptx.py"))
    )
    # tidak exec full (ada side-effect save), jadi fallback manual
    raise ImportError("skip import to avoid side-effect")
except Exception:
    pass

# Manual 10 slides (copy dari generate_pptx.py)
light_slides = [
    dict(
        title="Modul Performa Backend GR Demo",
        footer="Cover",
        cover=True,
        cover_title="Modul Performa\nBackend GR Demo",
        cover_sub="Logging + Performa  •  TIGA INSAN  •  60 menit",
        cover_bullets=[
            "Poster 1+2 + 10 Bab + 16 Endpoint + Checklist",
            "Postgres + Redis + pg_trgm + MatView + PgBouncer  →  p50 <50ms",
            "CDC Debezium WAL → Kafka → ES geospasial + ClickHouse",
        ],
    ),
    dict(
        title="Daftar Isi  —  10 Bab + 35 Slides Map",
        footer="Daftar Isi",
        bullets=[
            "Bab 1  Kecepatan = Kepercayaan (Muttaqin, SLA)  →  Slides 3-5",
            "Bab 2  Data Flow Flutter → Gateway → 7 Fondasi → 6 DB  →  Slide 13",
            "Bab 3  Postgres Scale (B-Tree, pg_trgm, MatView, Cursor, RLS)  →  Slides 14-17",
            "Bab 4  Caching (L1/L2/L3, TTL, tiering)  →  Slide 18  •  Bab 5 ES geospasial  →  Slide 19",
            "Bab 6  CDC Streaming (WAL → Kafka → ES/CH, anti dual-write)  →  Slide 20",
            "Bab 7  API Delivery (GZIP, cursor, Edge, rate limit)  →  Slide 21  •  Bab 8 Observability  →  Slide 22",
            "Bab 9  Roadmap 5 Fase (MVP Rp0 → auto-scale)  →  Slide 23  •  Bab 10 SLA 16 endpoint + throughput + DoD  →  Slides 24-26",
            "Demo 01-05 + Appendix A/B/C + Penutup  →  Slides 27-35",
        ],
    ),
    dict(
        title="Poster #1  —  Apa itu 200ms",
        footer="Poster #1",
        bullets=[
            "1 detik = 1000ms  —  200ms = 0.2 detik = kedipan mata.",
            "Diagram: User (tap) → Server (proses 50ms) → DB (query 20ms) → Response (GZIP 10ms) → User lihat hasil.",
            "200ms adalah batas 'terasa instan' — di atas itu pengguna mulai merasa menunggu (UX #46).",
            "Target backend <200ms p95 → sisa budget untuk jaringan 3G (500-1000ms) masih <3 detik total.",
            "Ukur dengan X-Response-Time header + prom-client histogram, bukan stopwatch manual.",
        ],
        mono="User --tap--> [Gateway JWT 5ms] --> [Redis HIT 2ms / DB 20ms] --> [GZIP 5ms] --> User\nTotal server: 30-50ms (p50)  |  + RTT 3G 500ms = ~550ms (masih <3s)",
    ),
    dict(
        title="Poster #3  —  Apa itu P99",
        footer="Poster #3",
        bullets=[
            "100 request diurutkan dari tercepat ke terlambat — P99 = request ke-99 (99% lebih cepat, 1% lebih lambat).",
            "Contoh: P99 = 200ms artinya 99 request <=200ms, 1 request >200ms (yang paling lambat, paling sering komplain).",
            "Kenapa peduli P99? Karena 1% yang lambat adalah pengguna yang paling vokal — mereka yang churn.",
            "Jangan hanya lihat rata-rata (avg) — avg bisa 50ms tapi P99 2000ms = ada yang sangat menderita.",
            "Spec: p50 <50ms, p95 <200ms, p99 <500ms — ketiganya harus hijau, bukan cuma p50.",
        ],
        mono="100 request sorted: [10,12,15,...,45,48, 50(p50), ..., 180, 200(p99), 2500]\n  50% <= p50   95% <= p95   99% <= p99   1% tail latency (paling lambat)",
    ),
    dict(
        title="Bab 3.1  —  Index B-Tree: 50.000x lebih cepat",
        footer="Bab 3",
        bullets=[
            "Tanpa index: Seq Scan baca 1.000.000 baris satu-per-satu — lambat, O(n).",
            "Dengan B-Tree: Balanced Tree cari dalam ~20 langkah (logaritmik) — O(log n) → 50.000x speedup.",
            "Aturan: semua FK wajib punya index (Checklist #3) — tanpa index, JOIN = Seq Scan.",
            "Buat CONCURRENTLY agar tidak blokir write di produksi; cek EXPLAIN: Index Scan vs Seq Scan.",
            "Index komposit: (community_id, pinned, created_at) untuk ORDER BY pinned, created_at LIMIT 20 (p50 <30ms).",
        ],
        mono="Tanpa index: [1M rows] --scan--> 1.000.000 langkah ~2000ms\nDengan B-Tree:  root -> branch -> branch -> leaf  ~20 langkah ~10ms\nCREATE INDEX CONCURRENTLY idx_komunitas_slug ON communities(slug);",
    ),
    dict(
        title="Bab 3.2  —  pg_trgm GIN: 10-50ms vs LIKE 2000ms",
        footer="Bab 3",
        bullets=[
            "LIKE '%ayam%' tidak bisa pakai B-Tree → Seq Scan 6.081 baris ~2000ms (meledak saat data tumbuh).",
            "pg_trgm pecah teks jadi trigram (3 huruf), GIN index cari overlap trigram → 10-50ms, similarity 71% untuk 'ayam'.",
            "Query: WHERE name % 'ayam' ORDER BY similarity(name,'ayam') DESC LIMIT 20 (operator % = similarity >0.3).",
            "Aktifkan: CREATE EXTENSION pg_trgm; CREATE INDEX USING GIN (name gin_trgm_ops) — gratis, built-in Postgres.",
            "Kapan ES? pg_trgm cukup untuk puluhan ribu dokumen; ES untuk 500+ komunitas / ratusan ribu dokumen (Fase 3).",
        ],
        mono="BEFORE: SELECT * FROM umkm WHERE name LIKE '%ayam%'  --> Seq Scan 2000ms\nAFTER : SELECT * FROM umkm WHERE name % 'ayam' ORDER BY similarity DESC --> GIN 10-50ms",
    ),
    dict(
        title="Bab 4  —  Caching: Hierarki L1/L2/L3 + TTL + Tiering",
        footer="Bab 4",
        bullets=[
            "Hierarki: L1 memory sub-ms (hot), L2 Redis 1-5ms (warm), L3 Postgres 10-50ms (cold) — kejar L1/L2.",
            "Pola: Cache-Aside (lazy, baca dulu cache, miss baru DB) vs Write-Through (tulis cache+DB bersamaan).",
            "TTL: jadwal sholat 1 jam, profil 5 menit, komunitas 10 menit, cari populer 5 menit — jangan cache selamanya.",
            "Tiering: hot (sering diakses, TTL pendek), warm (kadang, TTL menengah), cold (jarang, tanpa cache).",
            "Kapan Redis wajib? p99 cache >5ms atau >500 writes/detik atau hot data >80% hit rate (threshold 80%).",
        ],
        table_headers=["Tier", "Latency", "Contoh", "TTL"],
        table_rows=[
            ["L1 memory", "sub-ms", "Feature flag, JWT cache", "60s"],
            ["L2 Redis", "1-5ms", "Jadwal sholat, profil", "5m-1h"],
            ["L3 Postgres", "10-50ms", "Transaksi, ledger", "-"],
        ],
    ),
    dict(
        title="Bab 6  —  CDC: Debezium WAL → Kafka → ES/ClickHouse",
        footer="Bab 6",
        bullets=[
            "CDC = tangkap perubahan DB tanpa polling — Debezium baca WAL logical (INSERT/UPDATE/DELETE) → kirim ke Kafka.",
            "Kafka sebagai broker: consumer terpisah tulis ke ES (search) & ClickHouse (OLAP) — buffer & replay.",
            "Batch ETL alternatif untuk non real-time: cron 00:00 ekstrak Postgres → ClickHouse (laporan bulanan, hemat biaya).",
            "Anti dual-write: tulis HANYA ke Postgres (source of truth); jangan insert ke Postgres lalu ke ES di kode aplikasi — jika salah satu gagal = inkonsisten.",
            "Prinsip single writer: Postgres WAL adalah kebenaran; ES/ClickHouse adalah derived (bisa rebuild kapan saja).",
        ],
        mono="Postgres WAL --(Debezium pgoutput)--> Kafka(9092) --+--> ES(9200) geo <10ms\n                                           +--> ClickHouse(8123) OLAP\nDual-write DANGER: app->PG OK + app->ES FAIL = data hilang di search",
    ),
    dict(
        title="Bab 10.1  —  SLA 16 Endpoint (p50 / p95 / p99)",
        footer="Bab 10",
        bullets=[
            "16 endpoint dengan SLA per-endpoint — tidak ada satu angka untuk semua; baca tabel sebelum rilis.",
            "Read tercepat: /jadwal-sholat p50 <20ms (Redis TTL 1 jam); Write: /kas POST p50 <100ms (trigger SHA-256).",
            "Semua endpoint wajib EXPLAIN ANALYZE sebelum rilis; jika p95 >500ms → incident.",
            "Strategi utama per endpoint: cache, index, MatView, cursor, GIN — lihat kolom Strategi di tabel.",
        ],
        table_headers=["#", "Endpoint", "p50", "p95", "p99"],
        table_rows=[
            ["1", "GET /komunitas/:id", "<30ms", "<100ms", "<200ms"],
            ["3", "GET /kas", "<50ms", "<200ms", "<500ms"],
            ["5", "GET /pengumuman", "<30ms", "<100ms", "<200ms"],
            ["11", "GET /jadwal-sholat", "<20ms", "<50ms", "<100ms"],
            ["12", "GET /cari", "<50ms", "<200ms", "<500ms"],
            ["...", "16 endpoint total", "<50 read", "<200", "<500"],
        ],
    ),
    dict(
        title="Penutup  —  Kecepatan = Amanah",
        footer="Penutup",
        cover=True,
        cover_title="Kecepatan = Amanah",
        cover_sub="Mulai sederhana, tingkatkan bertahap — setiap ms = kepercayaan",
        cover_bullets=[
            "MVP Rp0 sudah p50 <50ms dengan Postgres+Redis+pg_trgm+MatView",
            "Tambah ES+CDC+ClickHouse saat 500+ komunitas (Fase 3) — jangan over-engineering di awal",
            "Ukur, jangan tebak: EXPLAIN ANALYZE, pg_stat_statements, Prometheus, Loki, Jaeger",
            "TIGA INSAN: Muttaqin (percaya), Shalih (berkarya), Nafi' (bermanfaat) — performa adalah wujudnya",
        ],
    ),
]

for idx, data in enumerate(light_slides, start=1):
    slide = prs.slides.add_slide(BLANK)
    if data.get("cover"):
        render_cover(slide, data, idx, total=len(light_slides))
        continue
    add_bg(slide)
    add_header(
        slide, data["title"], f"{data.get('footer', '')}  ·  {idx}/{len(light_slides)}"
    )
    bullets = data.get("bullets", [])
    has_table = "table_headers" in data
    has_mono = "mono" in data
    if has_table or has_mono:
        add_bullets(
            slide,
            bullets,
            left=Inches(0.4),
            top=Inches(1.25),
            width=Inches(6.9),
            height=Inches(4.0),
            font_size=Pt(9.5),
        )
        if has_table:
            add_table_slide(
                slide,
                data["table_headers"],
                data["table_rows"],
                Inches(7.6),
                Inches(1.25),
                Inches(5.3),
                Inches(4.2),
            )
        if has_mono:
            add_mono_box(
                slide,
                data["mono"],
                Inches(7.6),
                Inches(1.25) + (Inches(4.4) if has_table else Inches(0)),
                Inches(5.3),
                Inches(1.6) if has_table else Inches(4.2),
            )
    else:
        add_bullets(
            slide,
            bullets,
            left=Inches(0.4),
            top=Inches(1.25),
            width=Inches(12.5) if len(bullets) > 5 else Inches(7.0),
            height=Inches(5.6),
            font_size=Pt(10) if len(bullets) > 5 else Pt(10.5),
        )
    add_footer(slide, data.get("footer", ""), idx, total=len(light_slides))

out = str(pathlib.Path(__file__).with_name("Modul_Performa_Backend_GR_Demo_Light.pptx"))
prs.save(out)
print(f"Saved {out} with {len(prs.slides)} slides")
